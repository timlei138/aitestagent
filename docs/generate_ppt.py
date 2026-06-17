"""生成技术分享 PPT: AI Agent 驱动 Android 自动化测试"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 颜色方案 ──
DARK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x40, 0x9E, 0xFF)
GREEN = RGBColor(0x67, 0xC2, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
LIGHT_BLUE = RGBColor(0xEC, 0xF5, 0xFF)
ORANGE = RGBColor(0xE6, 0xA2, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_lines, title_color=BLUE):
    """Add a card-style box with title and body text."""
    card = add_shape(slide, left, top, width, height, LIGHT_BLUE)
    card.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFF)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Microsoft YaHei"

    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(4)
    return card


# ═══════════════════════════════════════
# Slide 1: 封面
# ═══════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK)

# 装饰条
add_shape(slide1, 0, 2.8, 13.333, 0.06, BLUE, MSO_SHAPE.RECTANGLE)

add_text_box(slide1, 1.5, 1.0, 10, 1.2,
             "AI Agent 驱动 Android 自动化测试",
             font_size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide1, 1.5, 2.0, 10, 0.7,
             "基于 LangChain + LangGraph 的技术实践",
             font_size=24, color=BLUE, align=PP_ALIGN.CENTER)

add_text_box(slide1, 1.5, 3.2, 10, 1.0,
             "技术分享会",
             font_size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide1, 1.5, 5.5, 10, 0.5,
             "2026-06-17  |  AI 自动化测试平台",
             font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# Slide 2: 项目概览 — 痛点与方案
# ═══════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

# 标题
add_shape(slide2, 0, 0, 13.333, 1.0, DARK, MSO_SHAPE.RECTANGLE)
add_text_box(slide2, 0.8, 0.15, 10, 0.7, "项目概览：传统测试 vs AI Agent 测试",
             font_size=28, color=WHITE, bold=True)

# 左侧：传统方式
add_card(slide2, 0.8, 1.4, 3.6, 3.2,
         "❌ 传统方式",
         ["✍️ 手工编写测试脚本",
          "🔧 每个 UI 变化都要改代码",
          "📍 XPath/ID 定位脆弱",
          "⏰ 维护成本高",
          "📝 学习成本高 (Appium/UiAutomator)"],
         RGBColor(0xF5, 0x6C, 0x6C))

# 箭头
add_text_box(slide2, 4.6, 2.5, 0.8, 0.6, "→", font_size=36, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# 右侧：AI Agent 方式
add_card(slide2, 5.5, 1.4, 3.6, 3.2,
         "✅ AI Agent 方式",
         ["🗣️ 自然语言描述测试目标",
          "🤖 LLM 理解页面 + 自主决策",
          "🎯 语义匹配元素 (label/role)",
          "🔄 UI 变化自适应",
          "📊 自动报告 + 知识积累"],
         GREEN)

# 右侧：核心能力
add_card(slide2, 9.8, 1.4, 3.0, 5.2,
         "🚀 核心能力",
         ["📥 自然语言输入",
          "🧠 LLM 自动规划",
          "🔧 Agent 自主执行",
          "✅ 自动验证结果",
          "👤 人机协作确认",
          "💾 RAG 知识复用",
          "📱 实时投屏监控"],
         BLUE)

# 底部技术栈
add_shape(slide2, 0.8, 5.0, 12, 1.8, LIGHT_BG)
add_text_box(slide2, 1.0, 5.1, 11.5, 1.6,
             "技术栈  |  LLM: DeepSeek V4 Pro + 智谱 GLM-4.6V  |  编排: LangGraph + LangChain  |  "
             "设备: uiautomator2  |  数据: ChromaDB + SQLite  |  "
             "API: FastAPI + WebSocket  |  前端: Vue 3 + Element Plus  |  "
             "Embedding: BGE-large-zh-v1.5 (本地)",
             font_size=14, color=DARK)

# ═══════════════════════════════════════
# Slide 3: LangGraph 核心 — 状态图 + 关键概念
# ═══════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

add_shape(slide3, 0, 0, 13.333, 1.0, DARK, MSO_SHAPE.RECTANGLE)
add_text_box(slide3, 0.8, 0.15, 10, 0.7, "LangGraph 状态图 — 多节点编排引擎",
             font_size=28, color=WHITE, bold=True)

# 状态图流程 (水平)
nodes = [
    ("Planner", "LLM 生成\n测试目标", BLUE),
    ("PlanReview", "interrupt()\n人机确认", ORANGE),
    ("Agent", "Tool-calling\n自主执行", GREEN),
    ("Reporter", "结果统计\n持久化", RGBColor(0x90, 0x94, 0x99)),
]
node_w = 2.5
gap = 0.5
start_x = 0.8
y_node = 1.5
for i, (name, desc, color) in enumerate(nodes):
    x = start_x + i * (node_w + gap)
    shape = add_shape(slide3, x, y_node, node_w, 1.5, color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.CENTER

    # 箭头 (除了最后一个)
    if i < len(nodes) - 1:
        ax = x + node_w + 0.05
        add_text_box(slide3, ax, y_node + 0.4, 0.5, 0.6, "→", font_size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# Agent 循环标注
add_shape(slide3, start_x + 2 * (node_w + gap), y_node + 1.7, node_w, 0.6, RGBColor(0xEC, 0xF5, 0xFF))
add_text_box(slide3, start_x + 2 * (node_w + gap) + 0.1, y_node + 1.75, node_w - 0.2, 0.5,
             "↻ 最多 12 轮迭代", font_size=14, color=BLUE, align=PP_ALIGN.CENTER)

# 下半部分：关键概念表格
concepts = [
    ("LangGraph 概念", "本项目应用", "文件位置"),
    ("StateGraph", "Planner→PlanReview→Agent→Reporter 状态机", "graph.py:340"),
    ("interrupt()", "PlanReview 节点暂停等待用户确认计划", "graph.py:294"),
    ("Command", "节点返回值，支持 update + goto 路由", "graph.py:287"),
    ("MemorySaver", "Checkpointer: 保存状态，支持 resume 恢复", "graph.py:351"),
    ("ToolNode", "Agent 子图中执行工具调用的节点", "graph.py:96"),
    ("tools_condition", "判断 LLM 输出是 ToolCall 还是 Text", "graph.py:92"),
    ("Subgraph", "Agent 内部独立 StateGraph (llm↔tools 循环)", "graph.py:95-99"),
]

table_top = 3.8
row_h = 0.42
for i, (c1, c2, c3) in enumerate(concepts):
    y = table_top + i * row_h
    bg_color = DARK if i == 0 else (LIGHT_BG if i % 2 == 0 else WHITE)
    fc = WHITE if i == 0 else DARK
    add_shape(slide3, 0.8, y, 2.5, row_h, bg_color, MSO_SHAPE.RECTANGLE)
    add_text_box(slide3, 0.9, y + 0.05, 2.3, row_h - 0.1, c1, font_size=13, color=fc, bold=(i == 0))

    add_shape(slide3, 3.4, y, 6.5, row_h, bg_color, MSO_SHAPE.RECTANGLE)
    add_text_box(slide3, 3.5, y + 0.05, 6.3, row_h - 0.1, c2, font_size=13, color=fc)

    add_shape(slide3, 10.0, y, 2.8, row_h, bg_color, MSO_SHAPE.RECTANGLE)
    add_text_box(slide3, 10.1, y + 0.05, 2.6, row_h - 0.1, c3, font_size=12, color=fc)

# ═══════════════════════════════════════
# Slide 4: 工具链 + 执行流程
# ═══════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)

add_shape(slide4, 0, 0, 13.333, 1.0, DARK, MSO_SHAPE.RECTANGLE)
add_text_box(slide4, 0.8, 0.15, 10, 0.7, "LangChain 工具链 & 执行流程演示",
             font_size=28, color=WHITE, bold=True)

# 左侧：17 个工具
tools_data = [
    ("🔍 感知类", ["get_screen_info — UI树解析", "query_app_knowledge — RAG检索", "query_element_identity — 元素身份"]),
    ("🖱️ 操作类", ["click — 语义点击 (核心)", "scroll_find_and_click — 滚动查找", "launch_app / type_input / press_key", "navigate_to / swipe / scroll_panel"]),
    ("✅ 验证类", ["assert_page_contains / assert_element_exists", "check_page_health / detect_popup", "recover_from_anomaly — 异常恢复", "wait_seconds — 等待"]),
]

ty = 1.3
for title, items in tools_data:
    add_card(slide4, 0.5, ty, 5.5, 0.3 + len(items) * 0.35, title, items, BLUE)
    ty += 0.5 + len(items) * 0.35 + 0.2

# 右侧：执行流程 Demo
add_shape(slide4, 6.5, 1.3, 6.3, 5.5, LIGHT_BG)
add_text_box(slide4, 6.8, 1.4, 5.8, 0.5, "🎬 执行流程 Demo", font_size=20, color=DARK, bold=True)

demo_steps = [
    ("1️⃣ 用户输入", "打开Settings, 通用设置→日期和时间→时区→英国", GREEN),
    ("2️⃣ Planner (LLM)", "生成目标: 验证时区设置为英国\nverification: [时区显示为英国]", BLUE),
    ("3️⃣ PlanReview", "interrupt() 暂停 → 前端弹窗 → 用户确认", ORANGE),
    ("4️⃣ Agent 执行", "perceive → click(通用设置)\n→ click(日期和时间) → click(时区)\n→ scroll_find_and_click(英国)", GREEN),
    ("5️⃣ 验证 DONE", "get_screen_info → 时间从 18:28→11:29 ✅\nDONE: 英国时区设置成功", BLUE),
    ("6️⃣ 报告", "Reporter → SQLite\n前端: 3步骤 / 100%通过率", GRAY),
]

dy = 2.1
for step, desc, color in demo_steps:
    box = add_shape(slide4, 6.8, dy, 5.8, 0.7, WHITE)
    box.fill.fore_color.rgb = WHITE
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.05)
    parts = desc.split("\n")
    p = tf.paragraphs[0]
    p.text = step + "  " + parts[0]
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.font.name = "Microsoft YaHei"
    p.font.bold = True
    if len(parts) > 1:
        p2 = tf.add_paragraph()
        p2.text = "      " + parts[1]
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        p2.font.name = "Microsoft YaHei"
    dy += 0.85

# Save
output_path = "docs/AI_Agent_Android_Test_Tech_Sharing.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
